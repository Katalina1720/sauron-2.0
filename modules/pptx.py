from io import BytesIO
from PIL import Image
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from modules.img_plugin import pillow_to_txt


def extract_chart_data(chart):
    chart_type = chart.chart_type
    chart_data = f"\nChart: {chart_type}"

    if chart.has_legend:
        legend = chart.legend
        for idx, legend_entry in enumerate(legend.entries):
            label = legend_entry.text
            chart_data += f"\nLegend {idx + 1}: {label}"

    for series in chart.series:
        series_name = series.name
        chart_data += f"\nSeries: {series_name}"
        for point in series.points:
            x = point.x if hasattr(point, 'x') else ''
            y = point.y if hasattr(point, 'y') else ''
            value = point.value
            chart_data += f"\n  Point: X={x}, Y={y}, Value={value}"

    return chart_data


def extract_table_data(table):
    table_data = "\nTable: "
    for row in table.rows:
        row_data = "\n|"
        for cell in row.cells:
            row_data += f" {cell.text_frame.text} |"
        table_data += row_data
    return table_data


def extract_image_data(image):
    img = Image.open(BytesIO(image.blob))
    image_data = '\nImage: '
    image_data += pillow_to_txt(img)
    return image_data


def parse_placeholder(placeholder):
    data = ''
    if placeholder.has_text_frame:
        if placeholder.text != '':
            data += "\n" + placeholder.text
    if placeholder.has_table == MSO_SHAPE_TYPE.CHART:
        data += extract_chart_data(placeholder.chart)
    if placeholder.has_chart == MSO_SHAPE_TYPE.TABLE:
        data += extract_table_data(placeholder.table)
    if placeholder.placeholder_format == MSO_SHAPE_TYPE.PICTURE:
        data += extract_image_data(placeholder.image)
    return data


def call_neded(shape):
    data = ''
    if not shape.is_placeholder and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        data += extract_image_data(shape.image)
    elif shape.is_placeholder:
        data += parse_placeholder(shape)
    else:
        if shape.has_chart:
            data += extract_chart_data(shape.chart)
        if shape.has_table:
            data += extract_table_data(shape.table)
        if shape.has_text_frame:
            if shape.text_frame.text != "":
                data += "\n" + shape.text_frame.text
    return data


def pars_shapes(cont):
    data = ''
    for shape in cont.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            data += pars_shapes(shape)
        else:
            data += call_neded(shape)
    return data


def read_file(path):
    prs = pptx.Presentation(path)
    data = ''
    i = 0
    for slide in prs.slides:
        data += f"\n\nSlide {i}:"
        data += pars_shapes(slide)
        i += 1

    return data
